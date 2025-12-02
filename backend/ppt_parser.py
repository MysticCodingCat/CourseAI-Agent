"""
PPT 解析器與結構預測
解析 PowerPoint 文件，提取結構與內容
"""
from typing import List, Dict, Any, Optional
from pathlib import Path
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTParser:
    """
    PowerPoint 解析器

    功能：
    1. 解析 PPT 檔案
    2. 提取文字、圖片、結構
    3. 識別章節與主題
    4. 生成初步講義大綱
    """

    def __init__(self):
        self.slides_data = []
        self.structure = {}

    def parse_ppt(self, ppt_path: str) -> List[Dict[str, Any]]:
        """
        解析 PPT 檔案

        Args:
            ppt_path: PPT 檔案路徑

        Returns:
            投影片資料列表
        """
        print(f"[PPT Parser] 解析 PPT: {ppt_path}")

        prs = Presentation(ppt_path)
        slides_data = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": self._extract_title(slide),
                "content": self._extract_content(slide),
                "notes": self._extract_notes(slide),
                "images": self._extract_images(slide),
                "type": self._identify_slide_type(slide)
            }

            slides_data.append(slide_data)

        print(f"[PPT Parser] 解析完成 - 共 {len(slides_data)} 張投影片")
        self.slides_data = slides_data
        return slides_data

    def _extract_title(self, slide) -> str:
        """提取投影片標題"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return ""

    def _extract_content(self, slide) -> List[str]:
        """提取投影片內容（文字）"""
        content = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # 跳過標題
            if shape == slide.shapes.title:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    content.append(text)

        return content

    def _extract_notes(self, slide) -> str:
        """提取備註"""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        return ""

    def _extract_images(self, slide) -> List[Dict[str, Any]]:
        """提取圖片資訊"""
        images = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({
                    "type": "image",
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                })
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # 可能包含圖片的佔位符
                if shape.has_text_frame:
                    continue
                images.append({
                    "type": "placeholder",
                    "name": shape.name
                })

        return images

    def _identify_slide_type(self, slide) -> str:
        """
        識別投影片類型

        返回：
        - title: 標題頁
        - section: 章節頁
        - content: 內容頁
        - summary: 總結頁
        - blank: 空白頁
        """
        title = self._extract_title(slide)
        content = self._extract_content(slide)

        # 空白頁
        if not title and not content:
            return "blank"

        # 標題頁（只有標題，內容很少）
        if title and len(content) <= 1:
            # 檢查標題關鍵字
            if any(keyword in title.lower() for keyword in ["課程", "lecture", "presentation", "簡報"]):
                return "title"

        # 章節頁（標題大、內容少）
        if title and len(content) <= 3:
            if any(keyword in title for keyword in ["章", "節", "Chapter", "Section", "Part"]):
                return "section"

        # 總結頁
        if title and any(keyword in title.lower() for keyword in ["總結", "結論", "summary", "conclusion", "回顧", "review"]):
            return "summary"

        # 內容頁
        return "content"

    def predict_structure(self) -> Dict[str, Any]:
        """
        預測 PPT 的章節結構

        Returns:
            結構資訊
        """
        print("[PPT Parser] 預測結構...")

        structure = {
            "title": "",
            "chapters": []
        }

        current_chapter = None

        for slide_data in self.slides_data:
            slide_type = slide_data["type"]

            # 標題頁
            if slide_type == "title":
                structure["title"] = slide_data["title"]
                continue

            # 章節頁 - 開始新章節
            if slide_type == "section":
                if current_chapter:
                    structure["chapters"].append(current_chapter)

                current_chapter = {
                    "title": slide_data["title"],
                    "slides": []
                }
                continue

            # 內容頁或總結頁 - 加入當前章節
            if slide_type in ["content", "summary"]:
                if current_chapter is None:
                    # 如果還沒有章節，創建一個
                    current_chapter = {
                        "title": "主要內容",
                        "slides": []
                    }

                current_chapter["slides"].append({
                    "slide_number": slide_data["slide_number"],
                    "title": slide_data["title"],
                    "type": slide_type
                })

        # 最後一個章節
        if current_chapter:
            structure["chapters"].append(current_chapter)

        self.structure = structure

        print(f"[PPT Parser] 結構預測完成 - {len(structure['chapters'])} 個章節")
        return structure

    def extract_keywords(self, slide_data: Dict[str, Any]) -> List[str]:
        """
        從投影片中提取關鍵詞

        Args:
            slide_data: 投影片資料

        Returns:
            關鍵詞列表
        """
        keywords = []

        # 從標題提取
        title = slide_data.get("title", "")
        if title:
            # 簡單的關鍵詞提取：標題本身
            keywords.append(title)

        # 從內容提取（尋找粗體、大寫、特殊格式的詞）
        content = slide_data.get("content", [])
        for line in content:
            # 簡單啟發式：
            # 1. 短句（可能是重點）
            if len(line) < 50:
                keywords.append(line)

            # 2. 包含冒號（可能是定義）
            if "：" in line or ":" in line:
                parts = line.split("：" if "：" in line else ":")
                if len(parts) == 2:
                    keywords.append(parts[0].strip())

        return keywords[:10]  # 最多 10 個關鍵詞

    def generate_outline_markdown(self) -> str:
        """
        生成 Markdown 講義大綱

        Returns:
            Markdown 文字
        """
        print("[PPT Parser] 生成 Markdown 大綱...")

        if not self.structure:
            self.predict_structure()

        md_lines = []

        # 標題
        md_lines.append(f"# {self.structure.get('title', '課程講義')}\n")
        md_lines.append("---\n")

        # 各章節
        for chapter in self.structure["chapters"]:
            md_lines.append(f"\n## {chapter['title']}\n")

            # 各投影片
            for slide_info in chapter["slides"]:
                slide_num = slide_info["slide_number"]
                slide_data = self.slides_data[slide_num - 1]

                md_lines.append(f"\n### {slide_data['title']}\n")
                md_lines.append(f"*投影片 {slide_num}*\n")

                # 內容（條列）
                if slide_data["content"]:
                    md_lines.append("\n**重點：**\n")
                    for line in slide_data["content"][:5]:  # 最多 5 條
                        md_lines.append(f"- {line}\n")

                # 關鍵詞
                keywords = self.extract_keywords(slide_data)
                if keywords:
                    md_lines.append(f"\n**關鍵詞：** {', '.join(keywords[:5])}\n")

                # 預留補充區域
                md_lines.append("\n**講義補充：**\n")
                md_lines.append("*（此區域將在課程進行時自動填充）*\n")

                md_lines.append("\n---\n")

        markdown_text = "".join(md_lines)

        print(f"[PPT Parser] Markdown 大綱生成完成 - {len(markdown_text)} 字元")
        return markdown_text

    def identify_topics_for_enhancement(self) -> List[Dict[str, Any]]:
        """
        識別需要網路搜尋補充的主題

        Returns:
            主題列表
        """
        print("[PPT Parser] 識別需要增強的主題...")

        topics = []

        for slide_data in self.slides_data:
            # 跳過非內容頁
            if slide_data["type"] not in ["content", "summary"]:
                continue

            keywords = self.extract_keywords(slide_data)

            for keyword in keywords:
                # 簡單啟發式：判斷是否需要補充
                needs_enhancement = self._check_if_needs_enhancement(keyword, slide_data)

                if needs_enhancement:
                    topics.append({
                        "keyword": keyword,
                        "slide_number": slide_data["slide_number"],
                        "slide_title": slide_data["title"],
                        "reason": needs_enhancement
                    })

        print(f"[PPT Parser] 識別出 {len(topics)} 個需要增強的主題")
        return topics

    def _check_if_needs_enhancement(self, keyword: str, slide_data: Dict[str, Any]) -> Optional[str]:
        """
        檢查關鍵詞是否需要補充

        Returns:
            需要補充的原因，或 None
        """
        # 啟發式規則

        # 1. 技術名詞（全英文）
        if keyword.isascii() and len(keyword) > 3:
            if keyword.isupper():  # 縮寫（例如 CNN、GAN）
                return "technical_term"
            if any(char.isupper() for char in keyword):  # 駝峰命名（例如 TensorFlow）
                return "technical_term"

        # 2. 包含特定關鍵字
        technical_keywords = ["演算法", "模型", "網路", "技術", "方法", "理論"]
        for tech_kw in technical_keywords:
            if tech_kw in keyword:
                return "technical_concept"

        # 3. 內容很少（可能需要補充）
        content_length = sum(len(line) for line in slide_data.get("content", []))
        if content_length < 100:
            return "insufficient_content"

        return None

    def save_to_json(self, output_path: str):
        """儲存解析結果到 JSON"""
        data = {
            "slides_data": self.slides_data,
            "structure": self.structure
        }

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        print(f"[PPT Parser] 儲存至 {output_path}")


# 測試函數
def test_ppt_parser():
    """測試 PPT 解析器"""
    print("=== 測試 PPT 解析器 ===\n")

    # 創建測試用的 PPT（需要手動準備）
    # 這裡假設有一個測試檔案

    parser = PPTParser()

    # 模擬解析結果（實際應該用真實 PPT）
    print("提示：請準備一個測試用的 PPT 檔案")
    print("然後使用：")
    print("  parser.parse_ppt('test.pptx')")
    print("  parser.predict_structure()")
    print("  markdown = parser.generate_outline_markdown()")
    print("  topics = parser.identify_topics_for_enhancement()")


if __name__ == "__main__":
    test_ppt_parser()
